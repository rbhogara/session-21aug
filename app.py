'''
Author(s) : Nikhil(nalampal@cisco.com) and Ravi(rbhogara@cisco.com)
Usage: Please refer the readme section of - https://github.com/rbhogara/talk-to-docs/

This script provides a Streamlit application for extracting text from PDF, DOCX, and PPTX files, generating embeddings using Ollama, storing it in ChromaDB, and querying the documents using Llama3 LLM.

Note: The purpose of the script is to have a starter code to build our own RAG applications with Local LLMs and Local Vector Database. This covers a basic use case of chatting with different types of documents
'''
import streamlit as st
import io
import PyPDF2
import docx
import pptx
from langchain_community.vectorstores import Chroma
from langchain_community import embeddings
from langchain_community.llms import Ollama
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain.text_splitter import CharacterTextSplitter
from collections import namedtuple

#Extracts text from a PDF file
def extract_text_from_pdf(uploaded_file):
    try:
        with io.BytesIO(uploaded_file.read()) as f:
            pdf_reader = PyPDF2.PdfReader(f)
            text = ""
            for page in range(len(pdf_reader.pages)):
                text += pdf_reader.pages[page].extract_text()
        return text
    except Exception as e:
        st.error(f"Error processing PDF file: {uploaded_file.name}")
        st.error(str(e))
        return None

#Extracts text from a DOCX file
def extract_text_from_docx(uploaded_file):
    try:
        with io.BytesIO(uploaded_file.read()) as f:
            doc = docx.Document(f)
            text = ""
            for para in doc.paragraphs:
                text += para.text
        return text
    except Exception as e:
        st.error(f"Error processing DOCX file: {uploaded_file.name}")
        st.error(str(e))
        return None

#Extracts text from a PPTX file
def extract_text_from_pptx(uploaded_file):
    try:
        with io.BytesIO(uploaded_file.read()) as f:
            presentation = pptx.Presentation(f)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text
        return text
    except Exception as e:
        st.error(f"Error processing PPTX file: {uploaded_file.name}")
        st.error(str(e))
        return None

#Processes the uploaded files to extract text and generate embeddings
def process_input(uploaded_files):
    try:
        docs_list = []
        for uploaded_file in uploaded_files:
            if uploaded_file.name.endswith('.pdf'):
                text = extract_text_from_pdf(uploaded_file)
            elif uploaded_file.name.endswith('.docx'):
                text = extract_text_from_docx(uploaded_file)
            elif uploaded_file.name.endswith('.pptx'):
                text = extract_text_from_pptx(uploaded_file)
            else:
                st.error(f"Unsupported file type: {uploaded_file.name}")
                continue
            if text:
                docs_list.append(text)

        #split data into chunks
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(" ".join(docs_list))

        #create Documents from chunks
        Document = namedtuple('Document', ['page_content', 'metadata'])
        documents = [Document(page_content=chunk, metadata={}) for chunk in chunks]

        #Generate embeddings with nomic-embed-text model and store them in a vectorstore
        vectorstore = Chroma.from_documents(
            documents=documents,
            collection_name="rag-chroma",
            embedding=embeddings.OllamaEmbeddings(model='nomic-embed-text'),
        )

        return vectorstore.as_retriever()
    except Exception as e:
        st.error("Error processing input files.")
        st.error(str(e))
        return None

#Answers a question based on the context retrieved from the documents
def answer_question(question, retriever):
    try:
        model_local = Ollama(model="llama3")
        after_rag_template = """Answer the question based only on the following context:
        {context}
        Question: {question}
        """
        after_rag_prompt = ChatPromptTemplate.from_template(after_rag_template)
        after_rag_chain = (
            {"context": retriever, "question": RunnablePassthrough()}
            | after_rag_prompt
            | model_local
            | StrOutputParser()
        )
        return after_rag_chain.invoke(question)
    except Exception as e:
        st.error("Error processing question.")
        st.error(str(e))
        return None

#Main function to run the Streamlit app
def main():
    st.title("Talk to Docs")
    st.write("Upload files and enter a question to query the documents.")

    uploaded_files = st.file_uploader("Upload files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)

    if uploaded_files:
        retriever = process_input(uploaded_files)
        if st.button('Generate Embeddings'):
            with st.spinner('Generating Embeddings...'):
                if retriever:
                    st.success('Embeddings generated successfully!')
                else:
                    st.error("Failed to generate embeddings.")

        question_input = st.text_input("Question")

        if st.button('Query Documents'):
            if not uploaded_files:
                st.warning("Please upload at least one file.")
            else:
                with st.spinner('Processing...'):
                    if retriever:
                        answer = answer_question(question_input, retriever)
                        st.text_area("Answer", value=answer, height=300, disabled=True)
                    else:
                        st.error("Failed to query documents.")

if __name__ == "__main__":
    main()
